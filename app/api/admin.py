from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Query,BackgroundTasks
from fastapi.responses import StreamingResponse, FileResponse
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.future import select
from sqlalchemy import delete, desc, func
from app.db.session import get_db, AsyncSessionLocal
from app.services.instruction_service import InstructionService
from app.services.eval_service import EvalService
from app.services.instruction_import_service import InstructionImportService
from app.services.rag_engine import RAGEngine
from app.models.base import Document, User, DocumentChunk, RAGTestRecord
from app.api.deps import get_current_user
from typing import List, Dict, Any, Optional
from pydantic import BaseModel
from app.core.route_logging import LoggingContextRoute
from app.core.config import settings
import uuid
import io
import httpx
import logging
import os
from datetime import datetime

# Import libraries for file processing
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pptx
except ImportError:
    pptx = None

router = APIRouter(route_class=LoggingContextRoute)
logger = logging.getLogger(__name__)

@router.post("/admin/instructions")
async def create_instruction(
    data: Dict[str, Any],
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    service = InstructionService(db)
    try:
        # Validate repository_id
        if 'repository_id' not in data:
            raise HTTPException(status_code=400, detail="repository_id is required")
        
        # Pass user_id to service
        return await service.create_instruction(data, user_id=current_user.id)
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@router.get("/admin/instructions")
async def list_instructions(
    page: int = 1,
    page_size: int = 10,
    repository_id: Optional[str] = None,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    service = InstructionService(db)
    repo_uuid = None
    if repository_id:
        try:
            repo_uuid = uuid.UUID(repository_id)
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid repository_id UUID")
            
    # Pass user_id to service
    return await service.get_instructions_paginated(user_id=current_user.id, page=page, page_size=page_size, repository_id=repo_uuid)

@router.post("/admin/instructions/import")
async def import_instructions(
    repository_id: str = Query(..., description="The repository ID to import instructions into"),
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    service = InstructionImportService(db)
    try:
        try:
            repo_uuid = uuid.UUID(repository_id)
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid repository_id UUID")

        content = await file.read()
        # Pass user_id, repository_id and filename to service
        return await service.import_from_excel(content, user_id=current_user.id, repository_id=repo_uuid, filename=file.filename)
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@router.post("/admin/eval/batch")
async def run_batch_eval(
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    service = EvalService(db)
    try:
        content = await file.read()
        # Pass user_id to service
        result_bytes = await service.run_batch_test(content, user_id=current_user.id)
        
        return StreamingResponse(
            io.BytesIO(result_bytes),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename=eval_result.xlsx"}
        )
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

from app.models.rag_config import RAGConfig, RAGConfigUpdate, RAGConfigResponse

@router.get("/admin/rag/config", response_model=RAGConfigResponse)
async def get_rag_config(
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    rag = RAGEngine(db)
    return await rag.get_config(current_user.id)

@router.put("/admin/rag/config", response_model=RAGConfigResponse)
async def update_rag_config(
    config_in: RAGConfigUpdate,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    rag = RAGEngine(db)
    config = await rag.get_config(current_user.id)
    
    for field, value in config_in.dict(exclude_unset=True).items():
        setattr(config, field, value)
    
    await db.commit()
    await db.refresh(config)
    return config

class RetrieveRequest(BaseModel):
    query: str
    doc_id: Optional[uuid.UUID] = None
    top_k: Optional[int] = 3

class RetrieveResultItem(BaseModel):
    id: uuid.UUID
    doc_id: uuid.UUID
    content: str
    score: float
    metadata: Dict[str, Any]

class RetrieveResponse(BaseModel):
    results: List[RetrieveResultItem]

@router.post("/admin/rag/retrieve", response_model=RetrieveResponse)
async def retrieve_test(
    request: RetrieveRequest,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    rag = RAGEngine(db)
    # Save test record
    record = RAGTestRecord(
        user_id=current_user.id,
        query=request.query,
        doc_id=request.doc_id
    )
    db.add(record)
    await db.commit()
    
    doc_ids = [request.doc_id] if request.doc_id else None
    
    try:
        results = await rag.search(
            query=request.query, 
            user_id=current_user.id, 
            top_k=request.top_k or 3, 
            doc_ids=doc_ids
        )
    except Exception as e:
        # Update record with error
        record.results = {"error": str(e)}
        await db.commit()
        raise HTTPException(status_code=500, detail=f"RAG Search failed: {str(e)}")
    
    # Update record with results
    # results is List[Dict], compatible with JSON
    record.results = results
    record.result_count = len(results)
    await db.commit()
    
    return {"results": results}

@router.get("/admin/rag/tests")
async def get_rag_tests(
    doc_id: Optional[uuid.UUID] = None,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    query = select(RAGTestRecord).filter(RAGTestRecord.user_id == current_user.id)
    
    if doc_id:
        query = query.filter(RAGTestRecord.doc_id == doc_id)
        
    stmt = query.order_by(RAGTestRecord.created_at.desc()).limit(50)
    result = await db.execute(stmt)
    return result.scalars().all()

class IndexRequest(BaseModel):
    provider: str
    model: str
    language: str = "zh"
    chunk_size: Optional[int] = None
    chunk_overlap: Optional[int] = None


async def background_index_document(doc_id: uuid.UUID, provider: str, model: str, chunk_size: Optional[int] = None, chunk_overlap: Optional[int] = None):
    """
    Background task for indexing documents.
    """
    async with AsyncSessionLocal() as session:
        try:
            logger.info(f"Background indexing started for doc_id: {doc_id}")
            
            # 1. Fetch Document
            stmt = select(Document).where(Document.id == doc_id)
            result = await session.execute(stmt)
            doc = result.scalar_one_or_none()
            
            if not doc:
                logger.error(f"Document {doc_id} not found in background task")
                return

            rag = RAGEngine(session)
            
            # 2. Clear existing chunks
            # We execute delete but don't commit yet. 
            # It will be committed together with new chunks in rag.index_document
            await session.execute(delete(DocumentChunk).where(DocumentChunk.doc_id == doc_id))
            
            # 3. Index Document (Generates embeddings and commits)
            # Note: rag.index_document commits the transaction!
            await rag.index_document(doc_id, doc.content, provider=provider, model=model, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
            
            # 4. Update Status (New Transaction)
            # Since expire_on_commit=False, 'doc' object is still valid but detached/clean.
            # We can modify it and commit.
            doc.status = "indexed"
            doc.is_configured = True
            doc.provider = provider
            doc.model = model
            doc.error_msg = None 
            if chunk_size:
                doc.chunk_size = chunk_size
            if chunk_overlap:
                doc.chunk_overlap = chunk_overlap
            
            await session.commit()
            logger.info(f"Background indexing finished for doc_id: {doc_id}")
            
        except Exception as e:
            logger.error(f"Background indexing failed for doc_id {doc_id}: {e}", exc_info=True)
            await session.rollback()
            
            # Re-fetch document to update status (in case of session issues)
            try:
                stmt = select(Document).where(Document.id == doc_id)
                result = await session.execute(stmt)
                doc = result.scalar_one_or_none()
                if doc:
                    doc.status = "failed"
                    doc.error_msg = str(e)
                    await session.commit()
            except Exception as ex:
                logger.error(f"Failed to update error status for doc_id {doc_id}: {ex}")

@router.post("/admin/documents/{doc_id}/index")
async def index_document_endpoint(
    doc_id: uuid.UUID,
    request: IndexRequest,
    background_tasks: BackgroundTasks,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    stmt = select(Document).where(Document.id == doc_id, Document.user_id == current_user.id)
    result = await db.execute(stmt)
    doc = result.scalar_one_or_none()
    
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
        
    # Set status to processing immediately
    doc.status = "processing"
    doc.language = request.language
    if request.chunk_size:
        doc.chunk_size = request.chunk_size
    if request.chunk_overlap:
        doc.chunk_overlap = request.chunk_overlap
    await db.commit()
    
    # Add background task
    background_tasks.add_task(background_index_document, doc_id, request.provider, request.model, request.chunk_size, request.chunk_overlap)
    
    return {"status": "processing", "message": "Indexing started in background"}

@router.get("/admin/ollama/models")
async def list_ollama_models(
    current_user: User = Depends(get_current_user)
):
    """
    Proxy request to Ollama to list available models.
    """
    if not settings.OLLAMA_API_BASE:
         raise HTTPException(status_code=400, detail="OLLAMA_API_BASE is not configured.")

    try:
        async with httpx.AsyncClient() as client:
            response = await client.get(f"{settings.OLLAMA_API_BASE}/api/tags")
            response.raise_for_status()
            data = response.json()
            # Extract model names
            models = [model["name"] for model in data.get("models", [])]
            return {"models": models}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch Ollama models: {str(e)}")

# Document Management Endpoints

@router.post("/admin/documents/upload")
async def upload_document(
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    try:
        content = ""
        filename = file.filename
        
        # Simple text extraction based on extension
        file_content = await file.read()
        
        # Save original file to uploads directory
        upload_dir = "uploads"
        if not os.path.exists(upload_dir):
            os.makedirs(upload_dir)
        saved_filename = f"{uuid.uuid4()}_{filename}"
        saved_path = os.path.join(upload_dir, saved_filename)
        with open(saved_path, "wb") as f:
            f.write(file_content)
        
        doc_status = "uploaded"
        error_msg = None

        if filename.lower().endswith(".pdf"):
            if pypdf:
                try:
                    pdf_reader = pypdf.PdfReader(io.BytesIO(file_content))
                    for page in pdf_reader.pages:
                        text = page.extract_text()
                        if text:
                            content += text + "\n"
                except Exception as e:
                    doc_status = "failed"
                    error_msg = f"Failed to parse PDF: {str(e)}"
                    logger.error(f"PDF parse error: {e}")
            else:
                doc_status = "failed"
                error_msg = "pypdf library not installed"
        elif filename.lower().endswith(".docx"):
            if docx:
                try:
                    doc_obj = docx.Document(io.BytesIO(file_content))
                    content = "\n".join([para.text for para in doc_obj.paragraphs])
                except Exception as e:
                    doc_status = "failed"
                    error_msg = f"Failed to parse DOCX: {str(e)}"
                    logger.error(f"DOCX parse error: {e}")
            else:
                doc_status = "failed"
                error_msg = "python-docx library not installed"
        elif filename.lower().endswith(".xlsx"):
            if openpyxl:
                try:
                    workbook = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        content += f"Sheet: {sheet_name}\n"
                        for row in sheet.iter_rows(values_only=True):
                            row_text = "\t".join([str(cell) for cell in row if cell is not None])
                            if row_text:
                                content += row_text + "\n"
                        content += "\n"
                except Exception as e:
                    doc_status = "failed"
                    error_msg = f"Failed to parse XLSX: {str(e)}"
                    logger.error(f"XLSX parse error: {e}")
            else:
                doc_status = "failed"
                error_msg = "openpyxl library not installed"
        elif filename.lower().endswith(".pptx"):
            if pptx:
                try:
                    prs = pptx.Presentation(io.BytesIO(file_content))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                        content += "\n"
                except Exception as e:
                    doc_status = "failed"
                    error_msg = f"Failed to parse PPTX: {str(e)}"
                    logger.error(f"PPTX parse error: {e}")
            else:
                doc_status = "failed"
                error_msg = "python-pptx library not installed"
        else:
            # Assume text/markdown/csv
            try:
                content = file_content.decode("utf-8")
            except UnicodeDecodeError:
                doc_status = "failed"
                error_msg = "File must be UTF-8 encoded text"

        # Sanitize content to remove null bytes which are not supported by PostgreSQL
        if content:
            content = content.replace("\x00", "")

        doc = Document(
            user_id=current_user.id,
            filename=filename,
            size=f"{len(file_content) / 1024:.2f} KB",
            status=doc_status,
            content=content,
            error_msg=error_msg,
            file_path=saved_path
        )
        db.add(doc)
        await db.commit()
        await db.refresh(doc)
        
        if doc_status == "failed":
            # If we want to notify frontend of failure but still keep the record
            # We can return the doc, but frontend needs to handle it.
            # Or we can raise exception, but then frontend receives error and refreshes list.
            # Since user wants to see the failure in the list, we should return successfully
            # but let the frontend see the status is 'failed'.
            pass
            
        return doc
    except Exception as e:
        logger.error(f"Upload failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

class DocumentResponse(BaseModel):
    id: uuid.UUID
    filename: str
    size: Optional[str]
    status: str
    created_at: datetime
    provider: Optional[str]
    model: Optional[str]
    language: Optional[str]
    chunk_size: Optional[int]
    chunk_overlap: Optional[int]
    is_configured: bool
    error_msg: Optional[str]
    
    class Config:
        orm_mode = True

class DocumentDetailResponse(DocumentResponse):
    content: Optional[str]

class PaginatedDocumentResponse(BaseModel):
    items: List[DocumentResponse]
    total: int
    page: int
    page_size: int
    total_pages: int

@router.get("/admin/documents/{doc_id}", response_model=DocumentDetailResponse)
async def get_document(
    doc_id: uuid.UUID,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    stmt = select(Document).where(Document.id == doc_id, Document.user_id == current_user.id)
    result = await db.execute(stmt)
    doc = result.scalar_one_or_none()
    
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
        
    return doc

@router.get("/admin/documents/{doc_id}/file")
async def get_document_file(
    doc_id: uuid.UUID,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    stmt = select(Document).where(Document.id == doc_id, Document.user_id == current_user.id)
    result = await db.execute(stmt)
    doc = result.scalar_one_or_none()
    
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
        
    if not doc.file_path or not os.path.exists(doc.file_path):
        raise HTTPException(status_code=404, detail="File not found on server")
        
    return FileResponse(doc.file_path, filename=doc.filename)

@router.get("/admin/documents", response_model=PaginatedDocumentResponse)
async def list_documents(
    page: int = Query(1, ge=1),
    page_size: int = Query(10, ge=1, le=100),
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    # Calculate offset
    offset = (page - 1) * page_size
    
    # Query total count
    count_stmt = select(func.count()).select_from(Document).where(Document.user_id == current_user.id)
    total_result = await db.execute(count_stmt)
    total = total_result.scalar()
    
    # Query documents
    # Use deferred to avoid loading content for the list view
    from sqlalchemy.orm import defer
    stmt = select(Document).where(Document.user_id == current_user.id)\
        .options(defer(Document.content))\
        .order_by(desc(Document.created_at))\
        .offset(offset).limit(page_size)
        
    result = await db.execute(stmt)
    documents = result.scalars().all()
    
    return {
        "items": documents,
        "total": total,
        "page": page,
        "page_size": page_size,
        "total_pages": (total + page_size - 1) // page_size
    }

@router.delete("/admin/documents/{doc_id}")
async def delete_document(
    doc_id: uuid.UUID,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    stmt = select(Document).where(Document.id == doc_id, Document.user_id == current_user.id)
    result = await db.execute(stmt)
    doc = result.scalar_one_or_none()
    
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
        
    # Delete associated document chunks
    await db.execute(delete(DocumentChunk).where(DocumentChunk.doc_id == doc_id))
    # Delete associated RAG test records
    await db.execute(delete(RAGTestRecord).where(RAGTestRecord.doc_id == doc_id))
    
    # Delete original file from server
    if doc.file_path and os.path.exists(doc.file_path):
        os.remove(doc.file_path)
        
    await db.delete(doc)
    await db.commit()
    return {"status": "success"}

@router.put("/admin/documents/{doc_id}/config_status")
async def update_document_config_status(
    doc_id: uuid.UUID,
    is_configured: bool,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    stmt = select(Document).where(Document.id == doc_id, Document.user_id == current_user.id)
    result = await db.execute(stmt)
    doc = result.scalar_one_or_none()
    
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
        
    doc.is_configured = is_configured
    await db.commit()
    return doc
